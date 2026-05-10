import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPM

def create_presentation():
    prs = Presentation()
    
    # Create a temp directory for PNGs
    if not os.path.exists("temp_pngs"):
        os.makedirs("temp_pngs")

    # Helper function to add a slide with title and content
    def add_slide(title, content=None, is_bullet=True):
        slide_layout = prs.slide_layouts[1] if content else prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title

        if content:
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            if isinstance(content, list):
                for item in content:
                    p = tf.add_paragraph()
                    p.text = item
                    p.level = 0
            else:
                tf.text = content
        return slide

    # Helper function to add an image slide
    def add_image_slide(title, img_path):
        slide_layout = prs.slide_layouts[5] # Title only
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        title_shape.text = title

        # Add image
        if os.path.exists(img_path):
            try:
                # Convert SVG to PNG if necessary
                if img_path.lower().endswith(".svg"):
                    png_path = os.path.join("temp_pngs", os.path.basename(img_path).replace(".svg", ".png"))
                    drawing = svg2rlg(img_path)
                    renderPM.drawToFile(drawing, png_path, fmt="PNG")
                    img_to_add = png_path
                else:
                    img_to_add = img_path

                left = Inches(1)
                top = Inches(1.5)
                width = Inches(8)
                slide.shapes.add_picture(img_to_add, left, top, width=width)
            except Exception as e:
                print(f"Error adding {img_path}: {e}")
                txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
                tf = txBox.text_frame
                p = tf.add_paragraph()
                p.text = f"Error adding image: {os.path.basename(img_path)}"
        else:
            print(f"File not found: {img_path}")

    # Slide 1: Title
    add_slide("Traffic Safety Mobile Application", "A unified platform for road safety and incident management")

    # Slide 2: Core Idea
    add_slide("Core Idea", [
        "Digitizes the full lifecycle of a traffic incident",
        "From reporting → police assessment → emergency dispatch → resolution",
        "Replaces slow phone-based reporting with a fast, GPS-based system"
    ])

    # Slide 3: Users & Roles - Driver
    add_slide("Users & Roles: Driver", [
        "Reports accidents/hazards with forms and media",
        "One-tap SOS with live GPS location",
        "Live traffic alerts and navigation"
    ])

    # Slide 4: Users & Roles - Police Officer
    add_slide("Users & Roles: Police Officer", [
        "Real-time incident alerts",
        "Review incident details (location, media, type)",
        "Dispatch emergency services or issue e-fines",
        "Send traffic alerts"
    ])

    # Slide 5: Users & Roles - Emergency Services
    add_slide("Users & Roles: Emergency Services", [
        "Receive dispatch with GPS location",
        "In-app routing to incident scene",
        "Request road clearance",
        "Status tracking (In Progress → Resolved)"
    ])

    # Slide 6: Users & Roles - Administrator
    add_slide("Users & Roles: Administrator", [
        "User management",
        "Configure alert rules",
        "Monitor analytics (incidents, response time, fines)",
        "Generate safety reports"
    ])

    # Slide 7: Key Features
    add_slide("Key Features", [
        "Real-time incident reporting (GPS + media)",
        "One-tap SOS alerts",
        "Live traffic incident map",
        "Role-based instant notifications",
        "Manual police severity decision",
        "In-app navigation with hazard overlays",
        "Electronic fines system",
        "Incident tracking & Admin analytics"
    ])

    # Slide 8: Technical Highlights
    add_slide("Technical Highlights", [
        "Cross-platform (Flutter / React Native)",
        "GPS accuracy up to ~10 meters",
        "Firebase push notifications (<10s delay)",
        "Secure authentication (JWT + bcrypt)",
        "HTTPS encrypted (TLS 1.2+)",
        "Cloud backend (500+ users support)",
        "Arabic + English (Full RTL support)"
    ])

    # Slide 9: Importance
    add_slide("Importance", [
        "Speeds up emergency response",
        "Improves coordination between authorities",
        "Structured, real-time data",
        "Turns smartphones into a safety network"
    ])

    # Diagrams
    diagrams = [
        ("Administrator Activity Diagram", r"out\Activity-Diagram\Admin\Administrator_Activity.svg"),
        ("Driver Activity Diagram", r"out\Activity-Diagram\Driver\Driver_Activity.svg"),
        ("Emergency Services Activity Diagram", r"out\Activity-Diagram\Emergency Services\EmergencyServices_Activity.svg"),
        ("Police Officer Activity Diagram", r"out\Activity-Diagram\PoliceOfficer\PoliceOfficer_Activity.svg"),
        ("Class Diagram", r"out\Class-Diagram\Class\Simplified_Traffic_Safety_ClassDiagram.svg"),
        ("Admin Use Case", r"out\plant-uml\Admin\Admin_UseCase.svg"),
        ("Driver Use Case", r"out\plant-uml\Driver\Driver_UseCase.svg"),
        ("Emergency Service Use Case", r"out\plant-uml\EmergencyService\Emergency_UseCase.svg"),
        ("Police Use Case", r"out\plant-uml\PoliceOfficer\Police_UseCase.svg"),
        ("Hacker Misuse Case", r"out\plant-uml-misuse-case\Hacker\Hacker_Misuse.svg"),
        ("Malicious Driver Misuse Case", r"out\plant-uml-misuse-case\Malicious Driver\MaliciousDriver_Misuse.svg"),
        ("Unauthorized User Misuse Case", r"out\plant-uml-misuse-case\Unauthorized User\UnauthorizedUser_Misuse.svg"),
        ("Report Incident Sequence Diagram", r"out\sequence Diagram\Report-incidents\Simplified_Report_Sequence.svg"),
        ("Assess & Dispatch Sequence Diagram", r"out\sequence Diagram\Simplified_Assess_Dispatch\Simplified_Assess_Dispatch_Sequence.svg")
    ]

    for title, path in diagrams:
        add_image_slide(title, path)

    # Save
    prs.save("Traffic_Safety_Mobile_App_Presentation.pptx")
    print("Presentation created successfully: Traffic_Safety_Mobile_App_Presentation.pptx")
    
    # Cleanup (optional)
    # import shutil
    # shutil.rmtree("temp_pngs")

if __name__ == "__main__":
    create_presentation()
